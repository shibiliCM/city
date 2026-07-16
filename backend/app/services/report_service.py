"""
Report Generation Service
--------------------------
Generates PDF (via WeasyPrint) and PPTX (via python-pptx) city intelligence reports.
All file bytes are stored in MongoDB GridFS and returned as download streams.
"""
from __future__ import annotations

import base64
import html
from io import BytesIO
from typing import Any

from motor.motor_asyncio import AsyncIOMotorDatabase, AsyncIOMotorGridFSBucket

from app.services.analytics_service import (
    get_accident_prone_areas,
    get_kpi_summary,
    get_pollution_hotspots,
    get_traffic_hotspots,
)
from app.utils.mongo import utcnow

try:
    from weasyprint import HTML as WeasyHTML
except Exception:
    WeasyHTML = None  # type: ignore[assignment]

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
except Exception:
    Presentation = None  # type: ignore[assignment]


# ─── HTML template helpers ────────────────────────────────────────────────────

def _kpi_cards_html(kpis: dict[str, Any]) -> str:
    icons = {
        "total_population": "👥",
        "avg_traffic_score": "🚗",
        "city_aqi": "🌫️",
        "accident_count": "⚠️",
        "city_health_score": "💚",
    }
    labels = {
        "total_population": "Total Population",
        "avg_traffic_score": "Avg Traffic Score",
        "city_aqi": "City AQI",
        "accident_count": "Accident Count",
        "city_health_score": "Health Score",
    }
    cards = ""
    for key, val in kpis.items():
        icon = icons.get(key, "📊")
        label = labels.get(key, key.replace("_", " ").title())
        fmt = f"{val:,.0f}" if isinstance(val, (int, float)) else str(val)
        cards += f"""
        <div class="kpi-card">
            <div class="kpi-icon">{icon}</div>
            <div class="kpi-value">{fmt}</div>
            <div class="kpi-label">{label}</div>
        </div>"""
    return cards


def _pdf_escape(value: str) -> str:
    return value.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")


def _simple_pdf(lines: list[str]) -> bytes:
    """Create a small valid PDF without native rendering dependencies."""
    wrapped: list[str] = []
    for raw in lines:
        text = str(raw).encode("latin-1", "replace").decode("latin-1")
        if not text:
            wrapped.append("")
            continue
        while len(text) > 92:
            split_at = text.rfind(" ", 0, 92)
            if split_at <= 0:
                split_at = 92
            wrapped.append(text[:split_at])
            text = text[split_at:].strip()
        wrapped.append(text)

    pages: list[list[str]] = []
    current: list[str] = []
    for line in wrapped:
        current.append(line)
        if len(current) >= 42:
            pages.append(current)
            current = []
    if current or not pages:
        pages.append(current)

    objects: list[bytes] = []
    page_refs: list[int] = []

    def add_object(payload: bytes) -> int:
        objects.append(payload)
        return len(objects)

    catalog_id = add_object(b"<< /Type /Catalog /Pages 2 0 R >>")
    pages_id = add_object(b"")
    font_id = add_object(b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>")

    for page in pages:
        content_lines = ["BT", "/F1 11 Tf", "50 780 Td", "14 TL"]
        for line in page:
            content_lines.append(f"({_pdf_escape(line)}) Tj")
            content_lines.append("T*")
        content_lines.append("ET")
        content = "\n".join(content_lines).encode("latin-1", "replace")
        content_id = add_object(b"<< /Length " + str(len(content)).encode() + b" >>\nstream\n" + content + b"\nendstream")
        page_id = add_object(
            f"<< /Type /Page /Parent {pages_id} 0 R /MediaBox [0 0 612 792] "
            f"/Resources << /Font << /F1 {font_id} 0 R >> >> /Contents {content_id} 0 R >>".encode()
        )
        page_refs.append(page_id)

    objects[pages_id - 1] = (
        f"<< /Type /Pages /Count {len(page_refs)} /Kids "
        f"[{' '.join(f'{ref} 0 R' for ref in page_refs)}] >>"
    ).encode()

    output = BytesIO()
    output.write(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets = [0]
    for index, payload in enumerate(objects, start=1):
        offsets.append(output.tell())
        output.write(f"{index} 0 obj\n".encode())
        output.write(payload)
        output.write(b"\nendobj\n")
    xref_offset = output.tell()
    output.write(f"xref\n0 {len(objects) + 1}\n".encode())
    output.write(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        output.write(f"{offset:010d} 00000 n \n".encode())
    output.write(
        f"trailer\n<< /Size {len(objects) + 1} /Root {catalog_id} 0 R >>\n"
        f"startxref\n{xref_offset}\n%%EOF\n".encode()
    )
    return output.getvalue()


def _fallback_pdf_bytes(
    city_id: str,
    report_type: str,
    kpis: dict[str, Any],
    traffic: list[dict[str, Any]],
    pollution: list[dict[str, Any]],
    accidents: list[dict[str, Any]],
) -> bytes:
    lines = [
        "CityTwin AI",
        f"{report_type.title()} Intelligence Report",
        f"City: {city_id}",
        f"Generated: {utcnow().strftime('%Y-%m-%d %H:%M UTC')}",
        "",
        "Executive Summary",
        f"Health Score: {float(kpis.get('city_health_score', 0)):.1f}/100",
        f"Total Population: {int(kpis.get('total_population', 0)):,}",
        f"Average Traffic Score: {float(kpis.get('avg_traffic_score', 0)):.1f}/100",
        f"City AQI: {float(kpis.get('city_aqi', 0)):.1f}",
        f"Accident Count: {int(kpis.get('accident_count', 0))}",
        "",
        "Top Traffic Hotspots",
    ]
    lines.extend(
        f"{i + 1}. {row.get('zone_name', row.get('zone_id', 'Unknown'))}: {float(row.get('traffic_score', 0)):.2f}"
        for i, row in enumerate(traffic[:10])
    )
    lines.extend(["", "Top Pollution Hotspots"])
    lines.extend(
        f"{i + 1}. {row.get('zone_name', row.get('zone_id', 'Unknown'))}: {float(row.get('aqi', 0)):.2f}"
        for i, row in enumerate(pollution[:10])
    )
    lines.extend(["", "Top Accident-Prone Areas"])
    lines.extend(
        f"{i + 1}. {row.get('zone_name', row.get('zone_id', 'Unknown'))}: {float(row.get('accident_density', 0)):.2f}"
        for i, row in enumerate(accidents[:10])
    )
    return _simple_pdf(lines)


def _hotspot_table_html(rows: list[dict], metric_key: str, metric_label: str) -> str:
    header = f"<tr><th>#</th><th>Zone</th><th>{metric_label}</th></tr>"
    body = "".join(
        f"<tr><td>{i + 1}</td><td>{r.get('zone_name', r.get('zone_id', 'Unknown'))}</td>"
        f"<td>{float(r.get(metric_key, 0)):.2f}</td></tr>"
        for i, r in enumerate(rows)
    )
    return f"<table><thead>{header}</thead><tbody>{body}</tbody></table>"


def _html_template(
    city_id: str,
    report_type: str,
    kpis: dict,
    traffic: list,
    pollution: list,
    accidents: list,
    recs: list,
    sims: list,
) -> str:
    generated_at = utcnow().strftime("%B %d, %Y at %H:%M UTC")
    kpi_cards = _kpi_cards_html(kpis)
    traffic_table = _hotspot_table_html(traffic, "traffic_score", "Traffic Score")
    pollution_table = _hotspot_table_html(pollution, "aqi", "AQI")
    accident_table = _hotspot_table_html(accidents, "accident_density", "Accident Density")

    rec_rows = "".join(
        f"<tr><td>{r.get('recommended_zone', r.get('zone_id', ''))}</td>"
        f"<td>{r.get('confidence_percent', r.get('confidence_score', 0))}</td>"
        f"<td>{r.get('query', '')[:80]}</td></tr>"
        for r in recs[:10]
    )
    sim_rows = "".join(
        f"<tr><td>{s.get('scenario_type', '')}</td>"
        f"<td>{s.get('delta_metrics', {}).get('traffic_change_pct', 0):.1f}%</td>"
        f"<td>{s.get('delta_metrics', {}).get('aqi_change_pct', 0):.1f}%</td>"
        f"<td>{s.get('confidence', 0):.0%}</td></tr>"
        for s in sims[:10]
    )

    return f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8" />
<title>CityTwin AI — {report_type.title()} Report</title>
<style>
  * {{ box-sizing: border-box; margin: 0; padding: 0; }}
  body {{ font-family: Inter, Arial, sans-serif; color: #1e293b; background: #f8fafc; font-size: 13px; }}
  .cover {{ background: linear-gradient(135deg, #0a0f1e 0%, #0d9488 100%);
             color: white; padding: 60px 50px; min-height: 200px; }}
  .cover h1 {{ font-size: 36px; font-weight: 700; margin-bottom: 8px; }}
  .cover .subtitle {{ font-size: 16px; opacity: 0.8; margin-bottom: 4px; }}
  .cover .meta {{ font-size: 12px; opacity: 0.6; margin-top: 12px; }}
  .section {{ padding: 32px 50px; }}
  .section + .section {{ border-top: 1px solid #e2e8f0; }}
  h2 {{ font-size: 20px; font-weight: 700; color: #0f172a; margin-bottom: 16px;
        padding-bottom: 8px; border-bottom: 2px solid #0d9488; }}
  h3 {{ font-size: 15px; font-weight: 600; color: #334155; margin: 16px 0 8px; }}
  .kpi-grid {{ display: flex; flex-wrap: wrap; gap: 12px; margin-bottom: 24px; }}
  .kpi-card {{ background: white; border: 1px solid #e2e8f0; border-radius: 12px;
               padding: 16px 20px; min-width: 140px; text-align: center;
               box-shadow: 0 1px 3px rgba(0,0,0,0.06); }}
  .kpi-icon {{ font-size: 24px; margin-bottom: 6px; }}
  .kpi-value {{ font-size: 22px; font-weight: 700; color: #0d9488; }}
  .kpi-label {{ font-size: 11px; color: #64748b; margin-top: 4px; }}
  table {{ width: 100%; border-collapse: collapse; margin-bottom: 20px;
           background: white; border-radius: 8px; overflow: hidden;
           box-shadow: 0 1px 3px rgba(0,0,0,0.06); }}
  th {{ background: #0f172a; color: white; padding: 10px 12px; text-align: left;
        font-size: 11px; text-transform: uppercase; letter-spacing: 0.05em; }}
  td {{ padding: 9px 12px; border-bottom: 1px solid #f1f5f9; }}
  tr:last-child td {{ border-bottom: none; }}
  tr:nth-child(even) td {{ background: #f8fafc; }}
  .exec-summary {{ background: #f0fdf4; border-left: 4px solid #0d9488;
                   padding: 16px 20px; border-radius: 0 8px 8px 0; margin-bottom: 20px; }}
  .footer {{ background: #0f172a; color: #64748b; padding: 20px 50px;
             font-size: 11px; text-align: center; }}
  @media print {{ body {{ background: white; }} }}
</style>
</head>
<body>

<div class="cover">
  <h1>🏙️ CityTwin AI</h1>
  <div class="subtitle">{html.escape(report_type.title())} Intelligence Report</div>
  <div class="subtitle">City: <strong>{html.escape(city_id)}</strong></div>
  <div class="meta">Generated {generated_at}</div>
</div>

<div class="section">
  <h2>Executive Summary</h2>
  <div class="exec-summary">
    <p>The city health score is <strong>{kpis.get('city_health_score', 0):.1f}/100</strong>
    with an average traffic congestion score of <strong>{kpis.get('avg_traffic_score', 0):.1f}</strong>
    and city-wide AQI of <strong>{kpis.get('city_aqi', 0):.1f}</strong>.
    Total population is <strong>{kpis.get('total_population', 0):,}</strong>
    with <strong>{kpis.get('accident_count', 0)}</strong> recorded accidents.
    </p>
  </div>
  <h2>City KPIs</h2>
  <div class="kpi-grid">{kpi_cards}</div>
</div>

<div class="section">
  <h2>Traffic Hotspots</h2>
  {traffic_table}
  <h2>Pollution Hotspots</h2>
  {pollution_table}
  <h2>Accident-Prone Areas</h2>
  {accident_table}
</div>

<div class="section">
  <h2>Planning Recommendations ({len(recs)} records)</h2>
  {"<table><thead><tr><th>Zone</th><th>Confidence</th><th>Query</th></tr></thead><tbody>" + rec_rows + "</tbody></table>" if recs else "<p>No recommendations yet. Run a planning query to generate recommendations.</p>"}
</div>

<div class="section">
  <h2>Simulation Results ({len(sims)} scenarios)</h2>
  {"<table><thead><tr><th>Scenario</th><th>Traffic Δ%</th><th>AQI Δ%</th><th>Confidence</th></tr></thead><tbody>" + sim_rows + "</tbody></table>" if sims else "<p>No simulations run yet.</p>"}
</div>

<div class="footer">
  CityTwin AI Platform &nbsp;|&nbsp; Confidential &nbsp;|&nbsp; {generated_at}
</div>

</body>
</html>"""


# ─── PDF generation ────────────────────────────────────────────────────────────

async def generate_pdf_report(
    db: AsyncIOMotorDatabase, city_id: str, report_type: str
) -> dict[str, Any]:
    kpis = await get_kpi_summary(db, city_id)
    traffic = await get_traffic_hotspots(db, city_id, 10)
    pollution = await get_pollution_hotspots(db, city_id, 10)
    accidents = await get_accident_prone_areas(db, city_id, 10)
    recs = await db.planning_recommendations.find(
        {"city_id": city_id}
    ).sort("created_at", -1).to_list(20)
    sims = await db.simulations.find(
        {"city_id": city_id}
    ).sort("created_at", -1).to_list(20)

    html = _html_template(
        city_id, report_type, kpis, traffic, pollution, accidents, recs, sims
    )

    if WeasyHTML is not None:
        try:
            pdf_bytes = WeasyHTML(string=html).write_pdf()
        except Exception:
            pdf_bytes = _fallback_pdf_bytes(city_id, report_type, kpis, traffic, pollution, accidents)
    else:
        pdf_bytes = _fallback_pdf_bytes(city_id, report_type, kpis, traffic, pollution, accidents)

    filename = f"{report_type}_{city_id}_{utcnow().strftime('%Y%m%d_%H%M%S')}.pdf"
    fs = AsyncIOMotorGridFSBucket(db)
    file_id = await fs.upload_from_stream(
        filename, pdf_bytes, metadata={"content_type": "application/pdf"}
    )
    return {
        "file_id": file_id,
        "content_type": "application/pdf",
        "filename": filename,
    }


# ─── PPTX generation ──────────────────────────────────────────────────────────

async def generate_pptx_report(
    db: AsyncIOMotorDatabase, city_id: str
) -> dict[str, Any]:
    if Presentation is None:
        # python-pptx not installed — return plain text stub
        stub = f"python-pptx is required. City: {city_id}".encode("utf-8")
        fs = AsyncIOMotorGridFSBucket(db)
        fid = await fs.upload_from_stream(
            f"citytwin_{city_id}.txt", stub,
            metadata={"content_type": "text/plain"}
        )
        return {"file_id": fid, "content_type": "text/plain", "filename": f"citytwin_{city_id}.txt"}

    kpis = await get_kpi_summary(db, city_id)
    traffic = await get_traffic_hotspots(db, city_id, 5)
    pollution = await get_pollution_hotspots(db, city_id, 5)
    recs = await db.planning_recommendations.find(
        {"city_id": city_id}
    ).sort("created_at", -1).to_list(5)
    sims = await db.simulations.find(
        {"city_id": city_id}
    ).sort("created_at", -1).to_list(5)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    DARK_BG = RGBColor(10, 15, 30)
    TEAL = RGBColor(13, 148, 136)
    WHITE = RGBColor(226, 232, 240)
    DIM = RGBColor(100, 116, 139)

    def add_slide(title: str, body_lines: list[str]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = DARK_BG

        # Title
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(1.1))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Teal underline bar
        from pptx.util import Emu
        bar = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0.7), Inches(1.55), Inches(12), Emu(60000)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = TEAL
        bar.line.fill.background()

        # Body
        body_tb = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.75), Inches(12), Inches(5.0)
        )
        body_tf = body_tb.text_frame
        body_tf.word_wrap = True
        for i, line in enumerate(body_lines[:14]):
            para = body_tf.paragraphs[0] if i == 0 else body_tf.add_paragraph()
            para.text = line
            para.font.size = Pt(16)
            para.font.color.rgb = DIM if line.startswith("  ") else WHITE

    # Slide 1: Title
    add_slide(
        f"🏙️ CityTwin AI — {city_id}",
        [
            "Urban Intelligence Report",
            f"Generated: {utcnow().strftime('%B %d, %Y')}",
            "",
            "Powered by Gemini 1.5 Pro · Prophet · XGBoost · NetworkX",
        ],
    )

    # Slide 2: City Health Overview
    add_slide(
        "City Health Overview",
        [
            f"  Health Score:     {kpis.get('city_health_score', 0):.1f}/100",
            f"  Population:       {kpis.get('total_population', 0):,}",
            f"  Avg Traffic Score:{kpis.get('avg_traffic_score', 0):.1f}",
            f"  City AQI:         {kpis.get('city_aqi', 0):.1f}",
            f"  Accident Count:   {kpis.get('accident_count', 0)}",
        ],
    )

    # Slide 3: Traffic Analysis
    add_slide(
        "Traffic Analysis",
        ["Top Congestion Zones:"] + [
            f"  {i+1}. {z.get('zone_name', z.get('zone_id',''))} — Score: {z.get('traffic_score',0):.1f}"
            for i, z in enumerate(traffic)
        ],
    )

    # Slide 4: Pollution Analysis
    add_slide(
        "Pollution Analysis",
        ["Top Pollution Zones (AQI):"] + [
            f"  {i+1}. {z.get('zone_name', z.get('zone_id',''))} — AQI: {z.get('aqi',0):.1f}"
            for i, z in enumerate(pollution)
        ],
    )

    # Slide 5: Forecasts
    add_slide(
        "Forecasts",
        [
            "Traffic: Prophet + XGBoost hybrid model",
            "Pollution: Prophet with seasonality regressors",
            "Population: Linear trend extrapolation to 2030",
            "Transport Demand: Random Forest on zone features",
            "",
            "Run forecast jobs from the Forecasting module",
            "to populate this slide with real projections.",
        ],
    )

    # Slide 6: Risk Assessment
    add_slide(
        "Risk Assessment",
        [
            "Risk types assessed per zone:",
            "  Flood Risk — Sigmoid model (elevation, rainfall, drainage)",
            "  Congestion Risk — Volume/capacity ratio",
            "  Pollution Risk — AQI forecast threshold mapping",
            "  Accident Risk — Logistic regression (history, road, signals)",
            "",
            "View interactive risk map in the CityTwin dashboard.",
        ],
    )

    # Slide 7: Planning Recommendations
    rec_lines = (
        [f"  {i+1}. {r.get('recommended_zone', r.get('zone_id',''))} — {r.get('query','')[:60]}"
         for i, r in enumerate(recs)]
        if recs else ["No recommendations yet — use the Planning module."]
    )
    add_slide("Planning Recommendations", ["Recent AI recommendations:"] + rec_lines)

    # Slide 8: Simulation Results
    sim_lines = (
        [f"  {i+1}. {s.get('scenario_type','')} — Traffic Δ: {s.get('delta_metrics',{}).get('traffic_change_pct',0):.1f}%"
         for i, s in enumerate(sims)]
        if sims else ["No simulations run yet — use the Simulation module."]
    )
    add_slide("Simulation Results", ["Recent scenario outcomes:"] + sim_lines)

    # Slide 9: Action Items
    add_slide(
        "Action Items",
        [
            "1. Validate field assumptions with on-ground surveys",
            "2. Prioritise zones with combined high traffic + AQI pressure",
            "3. Commission feasibility study for top planning recommendations",
            "4. Schedule simulation review with city engineering department",
            "5. Monitor KPIs weekly and re-run forecasts monthly",
            "6. Set up alert thresholds for AQI > 200 and traffic score > 80",
        ],
    )

    stream = BytesIO()
    prs.save(stream)
    data = stream.getvalue()
    filename = f"citytwin_{city_id}_{utcnow().strftime('%Y%m%d_%H%M%S')}.pptx"
    content_type = (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    fs = AsyncIOMotorGridFSBucket(db)
    file_id = await fs.upload_from_stream(
        filename, data, metadata={"content_type": content_type}
    )
    return {"file_id": file_id, "content_type": content_type, "filename": filename}
